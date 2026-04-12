"""
file_processor.py
Two responsibilities:
  1. extract_text(uploaded_file) -> (text, file_type_label)
  2. classify_document(text, filename) -> (accepted, label, reason)
     Uses the LLM to judge — no hardcoded keyword rules.
"""

import io
import json


# ── Text extraction ────────────────────────────────────────────────────────────

def extract_text(uploaded_file) -> tuple[str, str]:
    """Returns (text, detected_type_label)."""
    if uploaded_file is None:
        return "", ""

    filename = uploaded_file.name.lower()
    ext = filename.rsplit(".", 1)[-1] if "." in filename else ""
    raw = uploaded_file.read()

    extractors = {
        "pdf":  (_extract_pdf,   "PDF"),
        "docx": (_extract_docx,  "Word Document"),
        "doc":  (_extract_docx,  "Word Document"),
        "pptx": (_extract_pptx,  "PowerPoint"),
        "ppt":  (_extract_pptx,  "PowerPoint"),
        "xlsx": (_extract_xlsx,  "Excel Spreadsheet"),
        "xls":  (_extract_xlsx,  "Excel Spreadsheet"),
        "csv":  (_extract_plain, "CSV"),
        "txt":  (_extract_plain, "Text File"),
        "md":   (_extract_plain, "Markdown"),
        "json": (_extract_plain, "JSON"),
    }

    if ext in extractors:
        fn, label = extractors[ext]
        text = fn(raw)
    else:
        text = _extract_plain(raw)
        label = "Document"

    return text.strip(), label


def _extract_plain(raw: bytes) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return raw.decode(enc)
        except Exception:
            continue
    return ""


def _extract_pdf(raw: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)
    except Exception:
        return ""


def _extract_docx(raw: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw))
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    except Exception:
        return ""


def _extract_xlsx(raw: bytes) -> str:
    try:
        import pandas as pd
        dfs = pd.read_excel(io.BytesIO(raw), sheet_name=None)
        parts = []
        for name, df in dfs.items():
            parts.append(f"[Sheet: {name}]")
            parts.append(df.to_string(index=False))
        return "\n".join(parts)
    except Exception:
        return ""


# ── LLM-based document classifier ─────────────────────────────────────────────

_CLASSIFY_PROMPT = """You are a document classifier for a business consulting AI.

Decide if the following text is something a business consulting AI should analyse.

ACCEPT if the text is any of:
- A real business problem, challenge, or situation someone needs help solving
- A strategy document, business plan, pitch deck, or board update from a real company
- A workflow, process description, or operational issue
- A financial report with problems or gaps to address
- A case study that is being PRESENTED for actual problem-solving (not for academic teaching)

REJECT if the text is any of:
- A resume, CV, or personal profile
- A song, poem, lyrics, fiction, or creative writing
- An academic textbook, lecture notes, or theoretical essay (even if it contains business examples)
- A legal or contract document (terms, NDAs, clauses)
- Technical/software documentation (API docs, code, installation guides)
- A news article or blog post
- Personal diary, recipes, shopping lists, or unrelated personal content
- An academic assignment or homework question

Respond ONLY with a JSON object, no other text:
{{
  "accepted": true or false,
  "type": "short label for what this document is",
  "reason": "one direct sentence explaining the decision"
}}

Text to classify (first 1500 words):
\"\"\"
{text}
\"\"\"
"""


def classify_document(text: str, filename: str) -> tuple[bool, str, str]:
    """
    Uses the LLM to decide if this document is valid consulting input.
    Returns (accepted, type_label, reason).
    Falls back to accepting if the LLM call fails.
    """
    if not text or len(text.strip().split()) < 15:
        return False, "Empty / Unreadable", (
            "The file appears to be empty or could not be read. "
            "Check it isn't password-protected or an image-only scan."
        )

    # Send first ~1500 words to keep the classification call fast and cheap
    preview = " ".join(text.split()[:1500])

    result = _call_classifier(preview)
    if result is None:
        # LLM unavailable — fall back to accepting and let the main prompt judge
        return True, "Document", ""

    return result["accepted"], result.get("type", "Document"), result.get("reason", "")


def _call_classifier(text: str) -> dict | None:
    """Calls the fastest available LLM for a quick classification. Returns dict or None."""
    import os
    from dotenv import load_dotenv
    load_dotenv()

    prompt = _CLASSIFY_PROMPT.format(text=text)

    groq_key = os.getenv("GROQ_API_KEY", "")
    if groq_key:
        try:
            from groq import Groq
            client = Groq(api_key=groq_key)
            resp = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                max_tokens=120,
                temperature=0,
                messages=[{"role": "user", "content": prompt}],
            )
            raw = resp.choices[0].message.content.strip()
            return json.loads(raw)
        except Exception:
            pass

    anthropic_key = os.getenv("ANTHROPIC_API_KEY", "")
    if anthropic_key:
        try:
            import anthropic
            client = anthropic.Anthropic(api_key=anthropic_key)
            resp = client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=120,
                messages=[{"role": "user", "content": prompt}],
            )
            raw = resp.content[0].text.strip()
            return json.loads(raw)
        except Exception:
            pass

    return None
